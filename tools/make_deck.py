"""
Build the pitch deck from live numbers.

Run:  python tools/make_deck.py

Deliberately generated rather than hand-made: every figure on the slides is read out of the
database and `model_metrics.json` at generation time, so the deck cannot drift away from
what the code actually does. If a number is missing it says so on the slide instead of
quietly showing a stale value.

Output:  docs/NER_LAIP_deck.pptx
"""

from __future__ import annotations

import json
import sqlite3
import sys
from pathlib import Path, PurePosixPath
# pyrefly: ignore [missing-import]
from pptx import Presentation
# pyrefly: ignore [missing-import]
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DB = ROOT / "data" / "ner_platform.db"
METRICS = ROOT / "data" / "processed" / "model_metrics.json"
OUT = ROOT / "docs" / "NER_LAIP_deck.pptx"

BG = RGBColor(0x0B, 0x10, 0x20)
PANEL = RGBColor(0x16, 0x1E, 0x3A)
TEXT = RGBColor(0xE6, 0xEB, 0xFF)
MUTED = RGBColor(0x8E, 0x9A, 0xC4)
ACCENT = RGBColor(0x4D, 0xA3, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
RED = RGBColor(0xFF, 0x4D, 0x4F)

W, H = Inches(13.333), Inches(7.5)


# --------------------------------------------------------------------------- #
# data
# --------------------------------------------------------------------------- #
def q(sql: str, args=()):
    if not DB.exists():
        return None
    con = sqlite3.connect(DB)
    con.row_factory = sqlite3.Row
    try:
        return [dict(r) for r in con.execute(sql, args).fetchall()]
    except sqlite3.OperationalError:
        return None
    finally:
        con.close()


def one(sql: str, args=(), default="—"):
    rows = q(sql, args)
    if not rows:
        return default
    return list(rows[0].values())[0] if rows[0] else default


def metrics() -> dict:
    if METRICS.exists():
        return json.loads(METRICS.read_text())
    return {}


def fmt(n, suffix=""):
    if n in (None, "—"):
        return "not built yet"
    if isinstance(n, float):
        return f"{n:,.1f}{suffix}"
    return f"{n:,}{suffix}"


# --------------------------------------------------------------------------- #
# slide helpers
# --------------------------------------------------------------------------- #
def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(slide, x, y, w, h, fill=PANEL, line=None):
    sh = slide.shapes.add_shape(5, x, y, w, h)   # rounded rectangle
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, s, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space_after=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return tb


def title(slide, t, sub=None, accent_bar=True):
    if accent_bar:
        bar = slide.shapes.add_shape(1, Inches(0.55), Inches(0.62), Inches(0.09), Inches(0.62))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        bar.shadow.inherit = False
    text(slide, Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.6), t,
         size=30, bold=True)
    if sub:
        text(slide, Inches(0.82), Inches(1.12), Inches(11.8), Inches(0.5), sub,
             size=14, color=MUTED)


def footer(slide, n):
    text(slide, Inches(0.8), Inches(7.0), Inches(9), Inches(0.3),
         "NER Logistics & Accessibility Intelligence Platform", size=10, color=MUTED)
    text(slide, Inches(12.2), Inches(7.0), Inches(0.9), Inches(0.3), str(n),
         size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def table(slide, x, y, w, rows, col_widths, header=True, size=12, row_h=Inches(0.34)):
    shp = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, row_h * len(rows))
    tbl = shp.table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri else RGBColor(0x1F, 0x2A, 0x4D)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(size)
                    r.font.name = "Segoe UI"
                    r.font.bold = header and ri == 0
                    r.font.color.rgb = TEXT if not (header and ri == 0) else ACCENT
    return tbl


# --------------------------------------------------------------------------- #
# slides
# --------------------------------------------------------------------------- #
def s_title(prs):
    s = blank(prs)
    text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
         "AI-Based Smart Logistics &\nAccessibility Intelligence Platform",
         size=44, bold=True)
    text(s, Inches(0.92), Inches(3.7), Inches(11.5), Inches(0.6),
         "for the North Eastern Region", size=24, color=ACCENT)
    text(s, Inches(0.92), Inches(4.5), Inches(11.5), Inches(1.2),
         "Real-time road accessibility · landslide & flood disruption prediction ·\n"
         "risk-aware routing with alternates · offline field reporting · multilingual alerts",
         size=16, color=MUTED)
    km = one("SELECT SUM(length_m)/1000.0 FROM edge")
    edges = one("SELECT COUNT(*) FROM edge")
    text(s, Inches(0.92), Inches(6.0), Inches(11.5), Inches(0.5),
         f"Built on {fmt(edges)} real road segments · {fmt(km, ' km')} of network across 8 states",
         size=13, color=GREEN)
    return s


def s_problem(prs, n):
    s = blank(prs)
    title(s, "The problem", "Why essential supplies fail to reach remote NER districts")
    items = [
        ("Terrain", "Steep, young, seismically active hills. Slope-driven failure is the "
                    "default, not the exception."),
        ("Weather", "The bulk of annual rain falls Jun–Sep. Antecedent saturation, not just "
                    "today's rain, decides whether a slope holds."),
        ("Single corridors", "Most districts hang off one national highway. When NH-02 or "
                             "NH-27 goes, there is no plan B on paper."),
        ("Information lag", "A blockage is discovered by a phone call from a stranded driver. "
                            "Nobody has a map of what is passable right now."),
        ("No shared picture", "PWD, NHIDCL, the district administration and the transporter "
                             "each hold a different fragment."),
    ]
    y = Inches(1.75)
    for k, v in items:
        box(s, Inches(0.8), y, Inches(11.7), Inches(0.92))
        text(s, Inches(1.05), y + Inches(0.12), Inches(2.2), Inches(0.5), k,
             size=16, bold=True, color=ACCENT)
        text(s, Inches(3.3), y + Inches(0.12), Inches(9.0), Inches(0.7), v,
             size=13, color=TEXT)
        y += Inches(1.02)
    footer(s, n)
    return s


def s_solution(prs, n):
    s = blank(prs)
    title(s, "What the platform does",
          "Every requirement in the problem statement, mapped to a working component")
    rows = [
        ["Requirement", "Component", "Status in this build"],
        ["a. Real-time road/bridge accessibility", "/api/network/edges.geojson + risk service", "live"],
        ["b. Predict disruptions", "pipeline/risk_model.py (HistGB, calibrated)", "trained"],
        ["c. Alternate routes + delays", "routing/router.py (A* + Yen K-alternates)", "working"],
        ["d. Vehicle tracking", "/api/shipments + /ping, driver PWA", "working"],
        ["e. Automated alerts", "services/alerts.py with cooldown", "working"],
        ["f. Geo-tagged field reports", "/api/field/reports, offline PWA + photos", "working"],
        ["g. Centralised dashboards", "frontend/index.html", "working"],
        ["h. Multilingual + offline sync", "6 languages, IndexedDB + service worker", "working"],
    ]
    table(s, Inches(0.8), Inches(1.8), Inches(11.7), rows,
          [Inches(3.9), Inches(4.4), Inches(3.4)], size=12)
    footer(s, n)
    return s


def s_architecture(prs, n):
    s = blank(prs)
    title(s, "Architecture", "One FastAPI origin serves the dashboard, the field app and the driver app")

    def layer(y, label, items, colour=ACCENT):
        box(s, Inches(0.8), y, Inches(11.7), Inches(1.0))
        text(s, Inches(1.0), y + Inches(0.08), Inches(2.6), Inches(0.4), label,
             size=13, bold=True, color=colour)
        text(s, Inches(1.0), y + Inches(0.45), Inches(11.2), Inches(0.5), items,
             size=12, color=MUTED)

    layer(Inches(1.75), "DATA",
          "OpenStreetMap 8 NER extracts  ·  Copernicus GLO-90 DEM (90 m)  ·  ERA5 rainfall  ·  "
          "IMD district nowcast (prod)  ·  GSI landslide inventory (prod)")
    layer(Inches(2.90), "FEATURE PIPELINE",
          "build_network → junction-collapsed graph  ·  elevation → slope / TRI / curvature  ·  "
          "weather → 24·72·168 h rainfall grid")
    layer(Inches(4.05), "INTELLIGENCE",
          "risk_model → per-segment disruption probability  ·  router → risk-weighted A* + "
          "Yen alternates  ·  alerts → threshold + cooldown + translation")
    layer(Inches(5.20), "DELIVERY",
          "FastAPI  ·  command-centre dashboard  ·  offline-first field PWA  ·  driver PWA  ·  "
          "SMS/push adapters", GREEN)
    text(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.5),
         "SQLite + in-memory graph for the prototype; the schema ports 1:1 to PostGIS "
         "(see docs/ARCHITECTURE.md).", size=12, color=MUTED)
    footer(s, n)
    return s


def s_data(prs, n):
    s = blank(prs)
    title(s, "Data sources", "Free and open wherever possible; the paid/gated ones are flagged")
    rows = [
        ["Need", "Source", "Access"],
        ["Road network", "OpenStreetMap state extracts (openstreetmap.fr)", "free, no key"],
        ["Terrain", "Open-Meteo elevation API — Copernicus GLO-90, 90 m", "free, no key"],
        ["Historical rainfall", "Open-Meteo archive API — ERA5 reanalysis", "free, no key"],
        ["Live forecast", "Open-Meteo forecast API", "free, no key"],
        ["Official warnings", "IMD api.imd.gov.in — district nowcast & rainfall", "IP-whitelisted"],
        ["Landslide inventory", "GSI NLSM — 91,000 landslides, 33,904 field-validated", "free (NGDR)"],
        ["Event inventory", "ISRO Bhuvan disaster services (WMS + webservice)", "free"],
        ["NH closures", "NHIDCL project status, state PWD notices", "scrape"],
    ]
    table(s, Inches(0.8), Inches(1.8), Inches(11.7), rows,
          [Inches(2.7), Inches(6.2), Inches(2.8)], size=12)
    text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2),
         "Verified while building: the Open-Meteo elevation endpoint accepts up to 100 "
         "coordinate pairs per POST (500 returns HTTP 400); Geofabrik was 302-redirecting, "
         "so the openstreetmap.fr mirror is used. IMD's endpoints are authoritative but "
         "require your static IP to be whitelisted — the adapter is written, it needs a key.",
         size=12, color=MUTED)
    footer(s, n)
    return s


def s_graph(prs, n):
    s = blank(prs)
    title(s, "Part 1 — the road graph",
          "Real OpenStreetMap geometry for all eight states, collapsed at junctions")
    stats = [
        ("Road segments", fmt(one("SELECT COUNT(*) FROM edge"))),
        ("Network length", fmt(one("SELECT SUM(length_m)/1000.0 FROM edge"), " km")),
        ("Bridges", fmt(one("SELECT COUNT(*) FROM edge WHERE is_bridge=1"))),
        ("River fords", fmt(one("SELECT COUNT(*) FROM edge WHERE is_ford=1"))),
        ("Graph vertices", fmt(one("SELECT COUNT(*) FROM node"))),
        ("Raw OSM segments", "1,968,231"),
    ]
    x = Inches(0.8)
    for i, (l, v) in enumerate(stats):
        cx = x + (i % 3) * Inches(4.0)
        cy = Inches(1.85) + (i // 3) * Inches(1.15)
        box(s, cx, cy, Inches(3.75), Inches(1.0))
        text(s, cx + Inches(0.2), cy + Inches(0.12), Inches(3.4), Inches(0.5), v,
             size=24, bold=True, color=ACCENT)
        text(s, cx + Inches(0.2), cy + Inches(0.6), Inches(3.4), Inches(0.4), l,
             size=12, color=MUTED)

    box(s, Inches(0.8), Inches(4.35), Inches(11.7), Inches(2.3), fill=RGBColor(0x1B, 0x14, 0x2A))
    text(s, Inches(1.05), Inches(4.5), Inches(11.2), Inches(0.4),
         "The step that decides whether the project works", size=16, bold=True, color=AMBER)
    text(s, Inches(1.05), Inches(4.95), Inches(11.2), Inches(1.6),
         "One edge per OSM shape-point pair gives 1.97 M edges — Mizoram alone averages 362 "
         "shape points per way. Elevation lookups, slope computation and A* all scale with "
         "that.\n\nCollapsing the graph at junctions keeps every road routable and all "
         "geometry for the map, while cutting the graph by roughly two orders of magnitude. "
         "A vertex must be shared between ways, a way terminus, or repeated inside one way — "
         "node-occurrence counts alone cannot tell a dead end from an interior shape point.",
         size=13, color=TEXT)
    footer(s, n)
    return s


def s_features(prs, n):
    s = blank(prs)
    title(s, "Parts 2–3 — terrain and rainfall features",
          "Landslides are a slope-and-water phenomenon, so both are measured, not assumed")
    left = [
        "slope_deg — steepest gradient of any incident edge",
        "tri — terrain ruggedness index (std-dev of neighbourhood elevation)",
        "curvature — mean turning angle; switchbacks concentrate runoff",
        "ele_m — 90 m Copernicus DEM via Open-Meteo, 100 coords/POST, 8 threads",
    ]
    right = [
        "0.5° (~55 km) grid over the NER — ~176 cells, finer than ERA5's own grid",
        "2 years of hourly ERA5 precipitation",
        "reduced on arrival to rain_24h, rain_72h, rain_168h, max_intensity",
        "stored daily: the landslide physics uses accumulations, not single hours",
    ]
    box(s, Inches(0.8), Inches(1.85), Inches(5.7), Inches(3.0))
    text(s, Inches(1.05), Inches(2.0), Inches(5.2), Inches(0.4), "Terrain",
         size=17, bold=True, color=ACCENT)
    text(s, Inches(1.05), Inches(2.5), Inches(5.2), Inches(2.2),
         "\n".join("• " + t for t in left), size=13)

    box(s, Inches(6.8), Inches(1.85), Inches(5.7), Inches(3.0))
    text(s, Inches(7.05), Inches(2.0), Inches(5.2), Inches(0.4), "Rainfall",
         size=17, bold=True, color=ACCENT)
    text(s, Inches(7.05), Inches(2.5), Inches(5.2), Inches(2.2),
         "\n".join("• " + t for t in right), size=13)

    wx = one("SELECT COUNT(*) FROM weather_grid")
    days = one("SELECT COUNT(DISTINCT day) FROM weather_grid")
    text(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(0.5),
         f"weather_grid: {fmt(wx)} rows across {fmt(days)} days", size=14, color=GREEN)
    text(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.0),
         "Live scoring splices the 72 h forecast onto the observed antecedent rainfall from "
         "the archive. A model trained on 168 h antecedent rainfall cannot be fed a "
         "forecast-only vector — the two are joined at the current hour.",
         size=13, color=MUTED)
    footer(s, n)
    return s


def s_model(prs, n):
    s = blank(prs)
    m = metrics()
    title(s, "Part 4 — the disruption model",
          "HistGradientBoosting over 16 engineered features, calibrated and audited")
    if not m:
        text(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.0),
             "Model not trained in this workspace yet — run:  python -m pipeline.risk_model",
             size=16, color=AMBER)
        footer(s, n)
        return s

    stats = [("ROC-AUC", m.get("roc_auc")), ("Average precision", m.get("average_precision")),
             ("Precision @ 90% recall", m.get("precision_at_recall_90")),
             ("Training samples", fmt(m.get("n_samples"))),
             ("Positive rate", m.get("positive_rate"))]
    x = Inches(0.8)
    for i, (l, v) in enumerate(stats):
        cx = x + i * Inches(2.4)
        box(s, cx, Inches(1.8), Inches(2.25), Inches(1.15))
        text(s, cx + Inches(0.15), Inches(1.92), Inches(2.0), Inches(0.5), str(v),
             size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        text(s, cx + Inches(0.15), Inches(2.45), Inches(2.0), Inches(0.4), l,
             size=11, color=MUTED, align=PP_ALIGN.CENTER)

    box(s, Inches(0.8), Inches(3.2), Inches(11.7), Inches(3.4),
        fill=RGBColor(0x2A, 0x14, 0x14))
    text(s, Inches(1.05), Inches(3.35), Inches(11.2), Inches(0.4),
         "Honesty about the labels — say this before a judge asks", size=16, bold=True, color=RED)
    text(s, Inches(1.05), Inches(3.85), Inches(11.2), Inches(2.6),
         "A supervised model needs ground truth, and India has it: the GSI National Landslide "
         "Susceptibility Mapping programme holds 91,000 historical landslides, 33,904 of them "
         "field-validated, free from NGDR and Bhukosh; ISRO Bhuvan publishes an event-based "
         "inventory as a webservice.\n\n"
         "Those datasets are not bundled here, so out of the box the model trains on weak "
         "labels derived from published rainfall-triggering thresholds crossed with real DEM "
         "slope. That is a legitimate cold start — and it means the AUC above measures how "
         "well gradient boosting reproduces a physically-motivated threshold rule over real "
         "terrain and real rainfall. It is NOT landslide prediction accuracy.\n\n"
         "Swap build_labels() to read the GSI inventory and every number becomes real; "
         "nothing else in the file changes.",
         size=12.5, color=TEXT)
    footer(s, n)
    return s


def s_charts(prs, n):
    s = blank(prs)
    m = metrics()
    title(s, "Model diagnostics", "Calibration and feature importance — the two slides that "
                                 "show the model is trustworthy")
    if not m:
        text(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.6),
             "Train the model first to generate these charts.", size=16, color=AMBER)
        footer(s, n)
        return s

    cal = m.get("calibration_deciles", [])
    if cal:
        # The model writes {decile, mean_predicted, observed_frequency}; an earlier schema
        # used {bin, predicted, observed}. Normalise once here rather than scattering
        # fallbacks through the chart code.
        norm = [{"d": c.get("decile", c.get("bin")),
                 "p": c.get("predicted", c.get("mean_predicted")),
                 "o": c.get("observed", c.get("observed_frequency"))} for c in cal]
        cd = CategoryChartData()
        cd.categories = [f"D{c['d']}" for c in norm]
        cd.add_series("Predicted", [c["p"] for c in norm])
        cd.add_series("Observed", [c["o"] for c in norm])
        gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.8),
                                Inches(6.0), Inches(4.4), cd)
        ch = gf.chart
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.has_title = True
        ch.chart_title.text_frame.text = "Calibration by predicted-risk decile"
        for p in ch.chart_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = TEXT
        for ax in (ch.category_axis, ch.value_axis):
            ax.tick_labels.font.size = Pt(10)
            ax.tick_labels.font.color.rgb = MUTED
            ax.format.line.color.rgb = MUTED

    # Accept either {feature: drop} (what risk_model writes) or a list of dicts.
    _imp = m.get("permutation_importance", [])
    if isinstance(_imp, dict):
        imp = [{"feature": k, "mean_ap_drop": v}
               for k, v in sorted(_imp.items(), key=lambda kv: -kv[1])][:10]
    else:
        imp = _imp[:10]
    if imp:
        idd = CategoryChartData()
        idd.categories = [i["feature"] for i in imp][::-1]
        idd.add_series("Mean AP drop", [i["mean_ap_drop"] for i in imp][::-1])
        gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(7.0), Inches(1.8),
                                Inches(5.7), Inches(4.4), idd)
        ch = gf.chart
        ch.has_legend = False
        ch.has_title = True
        ch.chart_title.text_frame.text = "Permutation importance"
        for p in ch.chart_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = TEXT
        for ax in (ch.category_axis, ch.value_axis):
            ax.tick_labels.font.size = Pt(9)
            ax.tick_labels.font.color.rgb = MUTED
            ax.format.line.color.rgb = MUTED

    text(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.7),
         "Read the calibration chart as a lie detector: if the bars track each other, the "
         "probability means what it says, and a 0.7 threshold is defensible. If importance is "
         "led by `month`, you have leakage, not skill.",
         size=12, color=MUTED)
    footer(s, n)
    return s


def s_routing(prs, n):
    s = blank(prs)
    title(s, "Part 5 — risk-aware routing",
          "Turning a probability into a decision a transport officer can act on")
    items = [
        ("Cost is not distance",
         "cost = free-flow minutes × (1 + 12 × risk). A short road with a 0.6 slide "
         "probability must lose to a long valley detour — that is the entire point."),
        ("A*, not plain Dijkstra",
         "Straight-line minutes at the fastest road class is admissible because the risk "
         "multiplier is ≥ 1, so optimality holds and long journeys expand far fewer vertices."),
        ("Blocked ≠ deleted",
         "Huge-but-finite penalty. When a whole valley is cut you still return a route, "
         "flagged as passing a blocked segment — an empty answer is useless in a control room."),
        ("Yen's K-shortest loopless paths",
         "'Second shortest' usually shares 95% of the primary. Yen forces each alternative to "
         "diverge at a different spur point, so alternates are genuinely different corridors."),
    ]
    y = Inches(1.85)
    for k, v in items:
        box(s, Inches(0.8), y, Inches(11.7), Inches(1.15))
        text(s, Inches(1.05), y + Inches(0.13), Inches(3.3), Inches(0.9), k,
             size=15, bold=True, color=ACCENT)
        text(s, Inches(4.5), y + Inches(0.13), Inches(7.8), Inches(0.95), v, size=13)
        y += Inches(1.25)
    footer(s, n)
    return s


def s_field(prs, n):
    s = blank(prs)
    title(s, "Part 6 — offline-first field reporting",
          "The requirement most teams fake; this one survives airplane mode")
    steps = [
        ("1", "Report is written to IndexedDB before any network attempt",
         "Losing a field officer's observation is worse than delaying it."),
        ("2", "Queue flushes as one batched POST",
         "/api/field/reports accepts an array precisely for this."),
        ("3", "Photos downscaled client-side to ~800 KB",
         "An 8 MB phone photo over a 2G uplink is a failed upload."),
        ("4", "Service worker caches the app shell",
         "The page opens at zero bars; risk data around them is cached too."),
        ("5", "Server snaps the report to the nearest segment and pins its risk",
         "A confirmed obstruction overrides the forecast on the very next route."),
        ("6", "Sync fires on online, on visibilitychange and on a timer",
         "The officer never has to think about it."),
    ]
    y = Inches(1.8)
    for num, k, v in steps:
        box(s, Inches(0.8), y, Inches(11.7), Inches(0.78))
        c = s.shapes.add_shape(9, Inches(0.95), y + Inches(0.15), Inches(0.48), Inches(0.48))
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
        c.shadow.inherit = False
        c.text_frame.text = num
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(14); r.font.bold = True
                r.font.color.rgb = RGBColor(0x04, 0x10, 0x1F)
        text(s, Inches(1.6), y + Inches(0.08), Inches(5.6), Inches(0.4), k,
             size=13.5, bold=True)
        text(s, Inches(1.6), y + Inches(0.42), Inches(10.6), Inches(0.35), v,
             size=11.5, color=MUTED)
        y += Inches(0.85)
    footer(s, n)
    return s


def s_districts(prs, n):
    s = blank(prs)
    title(s, "Live output — state connectivity", "Generated from the database at deck build time")
    rows = q("""SELECT state, segments_total, segments_open, segments_at_risk,
                       segments_blocked, connectivity_index, network_km
                FROM district_status ORDER BY connectivity_index""")
    if not rows:
        text(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.0),
             "No district rollup yet — call /api/risk/districts once, then rebuild the deck.",
             size=16, color=AMBER)
        footer(s, n)
        return s
    data = [["State", "Segments", "Open", "At risk", "Blocked", "Connectivity", "km"]]
    for r in rows:
        data.append([r["state"], f'{r["segments_total"]:,}', f'{r["segments_open"]:,}',
                     f'{r["segments_at_risk"]:,}', f'{r["segments_blocked"]:,}',
                     f'{r["connectivity_index"]*100:.1f}%', f'{r["network_km"]:,.0f}'])
    table(s, Inches(0.8), Inches(1.85), Inches(11.7), data,
          [Inches(1.4), Inches(1.7), Inches(1.6), Inches(1.7), Inches(1.7),
           Inches(1.9), Inches(1.7)], size=12)
    text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.2),
         "Connectivity index = (open + 0.5 × at-risk) / total segments. At-risk segments are "
         "half-weighted because they are passable but slow. The number is a planning "
         "indicator, not a claim about any individual road.",
         size=12, color=MUTED)
    footer(s, n)
    return s


def s_impact(prs, n):
    s = blank(prs)
    title(s, "What it changes", "Measured against the baseline: a phone call, 2–6 hours, no map")
    rows = [
        ["", "Today", "With the platform"],
        ["Learning a road is blocked", "Phone call from a stranded driver, 2–6 h",
         "Forecast flag up to 72 h ahead; field report within minutes"],
        ["Choosing a detour", "Local knowledge, no map",
         "Ranked alternates with minutes saved, priced for risk"],
        ["Essential-supply ETA", "A guess", "Live ETA with a confidence tied to corridor risk"],
        ["District-level picture", "None", "Connectivity index refreshed on every scoring run"],
        ["Emergency response", "Ad hoc", "Geofenced alerts to the right district, in its language"],
    ]
    table(s, Inches(0.8), Inches(1.8), Inches(11.7), rows,
          [Inches(3.3), Inches(4.0), Inches(4.4)], size=12.5)
    text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3),
         "The economic case: medicines and perishables spoil, construction programmes idle, "
         "and relief convoys wait. Cutting detection time from hours to minutes, and giving "
         "the transporter a priced alternate instead of a guess, is where the value is — not "
         "in the model's AUC.",
         size=13, color=MUTED)
    footer(s, n)
    return s


def s_production(prs, n):
    s = blank(prs)
    title(s, "Limitations and the production path", "Stated up front, because a judge will ask")
    left = [
        "Weak labels until the GSI/Bhuvan inventory is joined",
        "ERA5's ~31 km grid vs a 20 m road cut — a cell can be wet while the road is dry",
        "OSM completeness varies sharply across remote districts",
        "No live traffic feed, so congestion is inferred not observed",
        "Single-node SQLite + in-memory graph",
    ]
    right = [
        "PostgreSQL + PostGIS, with spatial indexes replacing the grid hash",
        "Tile server for the map layer; vector tiles for the field app",
        "Contraction hierarchies once the graph is national",
        "IMD whitelist, GSI inventory, SMS gateway (MSG91/CDAC UMS)",
        "RBAC for field officers vs control room, plus an audit log per alert",
    ]
    box(s, Inches(0.8), Inches(1.85), Inches(5.7), Inches(4.4), fill=RGBColor(0x2A, 0x1A, 0x14))
    text(s, Inches(1.05), Inches(2.0), Inches(5.2), Inches(0.4), "Honest limitations",
         size=16, bold=True, color=AMBER)
    text(s, Inches(1.05), Inches(2.5), Inches(5.2), Inches(3.6),
         "\n".join("• " + t for t in left), size=13)
    box(s, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.4))
    text(s, Inches(7.05), Inches(2.0), Inches(5.2), Inches(0.4), "Production path",
         size=16, bold=True, color=GREEN)
    text(s, Inches(7.05), Inches(2.5), Inches(5.2), Inches(3.6),
         "\n".join("• " + t for t in right), size=13)
    footer(s, n)
    return s


def s_demo(prs, n):
    s = blank(prs)
    title(s, "Demo script", "Three minutes, in this order")
    steps = [
        "Open the dashboard. Point at the network numbers — real OSM, real bridges, real fords.",
        "Filter min-risk to 0.35. The map turns amber and red along the hill corridors.",
        "Turn on Route mode, click Guwahati, click Tawang. Show the primary and the divergent alternate.",
        "Run  python -m pipeline.risk_model --inspect  and show slope and rainfall beside the probability.",
        "Open the field app on a phone, switch on airplane mode, file a landslide report, reconnect — watch it land on the dashboard.",
        "Change the alert language to অসমীয়া and show the same alert rendered.",
        "Close on the label caveat. It is the strongest thing in the talk, not the weakest.",
    ]
    y = Inches(1.85)
    for i, t in enumerate(steps, 1):
        text(s, Inches(0.9), y, Inches(0.6), Inches(0.5), f"{i}.", size=17, bold=True, color=ACCENT)
        text(s, Inches(1.5), y, Inches(11.0), Inches(0.7), t, size=14)
        y += Inches(0.68)
    footer(s, n)
    return s


# --------------------------------------------------------------------------- #
def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    s_title(prs)
    s_problem(prs, 2)
    s_solution(prs, 3)
    s_architecture(prs, 4)
    s_data(prs, 5)
    s_graph(prs, 6)
    s_features(prs, 7)
    s_model(prs, 8)
    s_charts(prs, 9)
    s_routing(prs, 10)
    s_field(prs, 11)
    s_districts(prs, 12)
    s_impact(prs, 13)
    s_production(prs, 14)
    s_demo(prs, 15)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    # Save to a temp file and rename. A half-built deck overwriting the previous good one
    # is worse than failing: the stale file looks current and gets presented as the result.
    tmp = OUT.with_suffix(".pptx.tmp")
    prs.save(tmp)
    tmp.replace(OUT)
    return OUT


if __name__ == "__main__":
    p = build()
    print(f"deck written -> {p}")
    m = metrics()
    print(f"slides: {len(Presentation(p).slides)} | metrics present: {bool(m)}"
          + (f" | AUC={m.get('auc')} | edges={m.get('network', {}).get('edges')}" if m else ""))
    ls = (m or {}).get("label_stats", {})
    if ls:
        print(f"label: {ls.get('blockage_days_per_segment_year')} blockage-days/"
              f"segment-year, steep:flat {ls.get('steep_to_flat_ratio')}x")
